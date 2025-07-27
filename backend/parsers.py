# backend/parsers.py
import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation

def parse_pdf(contents: bytes) -> str:
    """Extracts text from a PDF file's bytes."""
    text = ""
    with fitz.open(stream=contents, filetype="pdf") as doc:
        for page in doc:
            text += page.get_text()
    return text

def parse_docx(contents: bytes) -> str:
    """Extracts text from a DOCX file's bytes."""
    text = ""
    doc = docx.Document(io.BytesIO(contents))
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def parse_pptx(contents: bytes) -> str:
    """Extracts text from a PPTX file's bytes."""
    text = ""
    prs = Presentation(io.BytesIO(contents))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text